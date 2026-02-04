import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_pitch_deck():
    prs = Presentation()

    # --- THEME COLORS ---
    # Based on "Holographic" app theme
    BG_COLOR = RGBColor(10, 15, 30)       # Dark Navy/Black
    ACCENT_COLOR = RGBColor(0, 229, 255)  # Cyan/Holographic Blue
    TEXT_COLOR = RGBColor(255, 255, 255)  # White
    SEC_COLOR = RGBColor(180, 180, 180)   # Light Grey

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title_slide(prs, title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_slide_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(54)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)

    def add_content_slide(prs, title_text, content_points):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.color.rgb = SEC_COLOR
            p.font.size = Pt(24)
            p.space_after = Pt(14)
            p.level = 0

    # --- SLIDE 1: TITLE ---
    add_title_slide(prs, "SecureTrack", "The Ultimate Anti-Theft & Asset Recovery Solution")

    # --- SLIDE 2: THE PROBLEM ---
    add_content_slide(prs, "The Problem", [
        "Smartphone theft is rising globally.",
        "Thieves immediately turn off devices to block tracking.",
        "Traditional 'Find My Device' fails when the phone is powered down.",
        "Data loss and privacy breaches cause massive stress."
    ])

    # --- SLIDE 3: THE SOLUTION ---
    add_content_slide(prs, "Enter SecureTrack", [
        "A military-grade protection suite for Android.",
        "Ensures your device NEVER truly sleeps.",
        "Works even when the thief thinks they've won.",
        "Stealth tracking via low-level SMS commands."
    ])

    # --- SLIDE 4: THE SECRET SAUCE ---
    add_content_slide(prs, "Core Technology: Fake Shutdown", [
        "When a thief tries to power off, we simulate a shutdown UI.",
        "Screen goes black, audio stops, vibration feedback mimics power-off.",
        "REALITY: The phone stays ON, tracking location and recording environment.",
        "The thief puts the 'dead' phone in their pocket, unknowingly broadcasting their location."
    ])

    # --- SLIDE 5: COMMAND CENTER ---
    add_content_slide(prs, "SMS Command Protocol", [
        "Control your device from ANY phone via SMS:",
        "#LOCATE_<PIN>: Triangulates & replies with Google Maps link.",
        "#SIREN_<PIN>: Blasts alarm at Max Volume (even if silent).",
        "#LOCK_<PIN>: Instantly locks the screen.",
        "#CALLME_<PIN>: Forces a callback to listen in."
    ])

    # --- SLIDE 6: SECURITY FEATURES ---
    add_content_slide(prs, "Defense & Recovery", [
        "Intruder Detection: Self-monitoring service.",
        "SIM Change Alert: Detects when SIM is swapped.",
        "Remote Wipe: Factory reset via SMS (Emergency Only).",
        "Holographic Dashboard: Professional, high-tech interface."
    ])

    # --- SLIDE 7: MARKET & BUSINESS ---
    add_content_slide(prs, "Market Opportunity", [
        "Target: High-value device owners, corporate fleets, parents.",
        "Business Model: Freemium.",
        "- Free: Basic Tracking & Siren.",
        "- Premium: Fake Shutdown & Stealth Recording.",
        "- Enterprise: Fleet management API."
    ])

    # --- SLIDE 8: CONCLUSION ---
    add_content_slide(prs, "Join the Revolution", [
        "SecureTrack isn't just an app; it's insurance.",
        "We turn the thief's advantage into their downfall.",
        "Available now on GitHub.",
        "\nContact: Iamjunade"
    ])

    output_file = "SecureTrack_Pitch.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_pitch_deck()
